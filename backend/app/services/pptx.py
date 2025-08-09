import io
import uuid
from pathlib import Path
from typing import Callable, List, Optional, Tuple

import requests
from pptx import Presentation
from pptx.presentation import Presentation as PresentationType
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

from app.core.config import settings
from app.models.slides import SlidePlan, ImageMeta

BRANDED_PLACEHOLDER = settings.STABILITY_PLACEHOLDER_URL


class TemplateError(Exception):
    pass


class ImagePlacementError(Exception):
    pass


def _load_template() -> PresentationType:
    try:
        if settings.PPTX_TEMPLATE_PATH:
            return Presentation(settings.PPTX_TEMPLATE_PATH)
        return Presentation()
    except Exception as e:
        raise TemplateError(f"Failed to load PPTX template: {str(e)}")


def _download_image(url: str) -> Optional[bytes]:
    try:
        resp = requests.get(url, timeout=settings.PPTX_IMAGE_HTTP_TIMEOUT_SECONDS)
        resp.raise_for_status()
        return resp.content
    except Exception:
        return None

# Backwards-compatible alias for tests expecting download_image symbol
def download_image(url: str) -> Optional[bytes]:
    return _download_image(url)


def _add_title(slide, title: str) -> None:
    slide.shapes.title.text = title
    tf = slide.shapes.title.text_frame
    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.name = settings.PPTX_FONT_NAME
            run.font.size = Pt(settings.PPTX_TITLE_FONT_SIZE_PT)
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT


def _add_bullets(slide, bullets: List[str]) -> None:
    # Assume placeholder index 1 is content
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        p = body.add_paragraph() if i > 0 else body.paragraphs[0]
        p.text = bullet
        p.level = 0
        p.font.name = settings.PPTX_FONT_NAME
        p.font.size = Pt(settings.PPTX_BODY_FONT_SIZE_PT)
        p.word_wrap = True


def _fit_size(img_width_in: float, img_height_in: float) -> Tuple[float, float]:
    max_w = settings.PPTX_IMAGE_MAX_WIDTH_IN
    max_h = settings.PPTX_IMAGE_MAX_HEIGHT_IN
    scale = min(max_w / img_width_in, max_h / img_height_in, 1.0)
    return img_width_in * scale, img_height_in * scale


def _add_image(slide, img_bytes: bytes) -> None:
    # Save into memory and insert; python-pptx needs a stream or file path
    stream = io.BytesIO(img_bytes)
    # Approximate size placement; library lacks pre-read size without PIL, so use max sizes directly
    width, height = _fit_size(settings.PPTX_IMAGE_MAX_WIDTH_IN, settings.PPTX_IMAGE_MAX_HEIGHT_IN)
    left = Inches((10 - width) / 2)  # assuming 10 inches slide width default
    top = Inches(2)
    slide.shapes.add_picture(stream, left, Inches(2), width=Inches(width), height=Inches(height))


def _set_image_alt_text(slide, alt_text: str) -> None:
    try:
        pic = slide.shapes[-1]
        if hasattr(pic, "alternative_text"):
            pic.alternative_text = alt_text
    except Exception:
        pass


def build_pptx(
    slides: List[SlidePlan],
    output_dir: str | None = None,
    on_progress: Optional[Callable[[int, int], None]] = None,
) -> Path:
    """
    Build a PPTX file from SlidePlan objects with images and notes.
    Emits progress via callback (completed, total) if provided.
    """
    prs = _load_template()
    total = len(slides)

    for idx, slide_plan in enumerate(slides, start=1):
        try:
            layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(layout)
            _add_title(slide, slide_plan.title)
            _add_bullets(slide, slide_plan.bullets)

            if slide_plan.image:
                img_bytes = _download_image(slide_plan.image.url)
                if not img_bytes:
                    img_bytes = _download_image(BRANDED_PLACEHOLDER)
                if img_bytes:
                    _add_image(slide, img_bytes)
                    _set_image_alt_text(slide, slide_plan.image.alt_text)

            if slide_plan.notes:
                slide.notes_slide.notes_text_frame.text = slide_plan.notes
        except Exception:
            # Continue on error per slide
            pass
        finally:
            if on_progress:
                try:
                    on_progress(idx, total)
                except Exception:
                    pass

    out_dir = Path(output_dir or settings.PPTX_TEMP_DIR)
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / f"{uuid.uuid4()}.pptx"
    prs.save(str(out_path))
    return out_path